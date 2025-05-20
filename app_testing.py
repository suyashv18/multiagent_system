import streamlit as st
from langgraph.graph import StateGraph, END
from langchain.chains import ConversationalRetrievalChain, LLMChain
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import PyPDFLoader
from langchain.memory import ConversationBufferMemory
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_groq import ChatGroq
from langchain.prompts import ChatPromptTemplate, PromptTemplate
from langchain_community.utilities import SQLDatabase
from typing import TypedDict, Optional
from langchain.sql_database import SQLDatabase
from pptx import Presentation
from pptx.util import Inches
import sys 
import sqlite3
import pandas as pd

# Predefined API key and vector store paths
api_key = st.secrets["GROQ_API_KEY"]
embedding = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
llm = ChatGroq(
    groq_api_key=api_key,
    model_name="llama-3.3-70b-versatile"
)

# Load pre-built vector stores (paths are predefined)
@st.cache_resource
def load_vectorstores():
    ml_vectorstore = FAISS.load_local("transformers_vdb", embedding, allow_dangerous_deserialization=True)
    medical_vectorstore = FAISS.load_local("radiotherapy_vdb", embedding, allow_dangerous_deserialization=True)
    sdlc_vectorstore = FAISS.load_local("sdlc_vdb", embedding, allow_dangerous_deserialization=True)
    # finance_vectorstore = FAISS.load_local("finance_vdb", embedding, allow_dangerous_deserialization=True)
    return ml_vectorstore, medical_vectorstore, sdlc_vectorstore

ml_vectorstore, medical_vectorstore, sdlc_vectorstore = load_vectorstores()

finance_prompt = PromptTemplate.from_template("""
You are a financial data extraction expert. Your job is to answer questions about financial data and ALWAYS present the answers in a clear markdown table format.

INSTRUCTIONS:
1. Extract relevant financial data from the provided context
2. Structure ALL responses as markdown tables with proper headers
3. Use numerical formatting for financial figures (e.g., currency, percentages)
4. Include column headers that clearly describe the data
5. Include a markdown table separator row with dashes after the header row

Context:
-----------------
{context}
-----------------

Chat History:
{chat_history}

Question: {question}

YOUR ANSWER MUST BE IN MARKDOWN TABLE FORMAT like this:
| Header1 | Header2 | Header3 |
| --- | --- | --- |
| Value1 | Value2 | Value3 |
| Value4 | Value5 | Value6 |

Answer: 
""")

# Create QA chains for each agent
def create_qa_chain(vectorstore):
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
    qa_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(search_kwargs={"k": 4}),
        memory= memory
    )
    return qa_chain

agent_a_chain = create_qa_chain(ml_vectorstore)
agent_b_chain = create_qa_chain(medical_vectorstore)
agent_c_chain = create_qa_chain(sdlc_vectorstore)

# Create the finance chain with custom prompt
memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
loader = PyPDFLoader("data/Finance.pdf")
pages = loader.load_and_split()
vectordb = FAISS.from_documents(pages, embedding)

agent_d_chain = ConversationalRetrievalChain.from_llm(
    llm =llm,
    retriever= vectordb.as_retriever(search_kwargs={"k": 4}),
    memory = memory,
    combine_docs_chain_kwargs={"prompt": finance_prompt}
)

def create_chain():
    # Load FAISS vector store
    vectorstore = FAISS.load_local("faiss_index", embeddings=embedding,allow_dangerous_deserialization=True )
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
    qa_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(search_kwargs={"k": 4}),
        memory=memory
    )
    print("QA Chain", qa_chain)
    return qa_chain

def extract_slide_content(text):
    prompt = f"""
    Extract key points from the following document that would be useful in a presentation.
    Output should be in the form:
    - Title Slide: [Title]
    - Slide 1: [Title] | [Bullet Points] Always return the bullet points with * in the beginning.
    - Slide 2: [Title] | [Bullet Points] [TABLE:] If any tabular data is present return it in tabular format with TABLE: in the beginning.
    ...

    Remember : For cost related details return data always in tabular format.
    Always return the details in the above format only.
    Document:
    {text}
    """

    prompt = PromptTemplate.from_template(prompt)
    formatted_prompt = prompt.format(text=text)
    response = llm.invoke(formatted_prompt)
    selected = response.content.strip()

    return selected

def parse_slide_content(text):
    lines = text.strip().split('\n')
    slides = []
    current_slide = None
    in_table = False
    table_rows = []

    for line in lines:
        line = line.strip()
        if line.startswith("- Title Slide:"):
            current_slide = {
                "title": line.replace("- Title Slide:", "").strip(),
                "bullets": [],
                "table": [],
                "is_title": True
            }
            slides.append(current_slide)

        elif line.startswith("- Slide"):
            if current_slide and in_table:
                current_slide["table"] = table_rows
                table_rows = []
                in_table = False

            parts = line.split(":")
            title_and_content = parts[1].split("|")
            title = title_and_content[0].strip()
            current_slide = {
                "title": title,
                "bullets": [],
                "table": [],
                "is_title": False
            }
            slides.append(current_slide)
        elif line.startswith("•"):
            current_slide["bullets"].append(line.replace("•", "").strip())

        elif line.startswith("*"):
            current_slide["bullets"].append(line.replace("*", "").strip())

        elif line.startswith("TABLE:"):
            in_table = True
            table_rows = []

        elif in_table and "|" in line:
            row = [cell.strip() for cell in line.strip("|").split("|")]
            table_rows.append(row)

    if current_slide and in_table:
        current_slide["table"] = table_rows

    create_ppt_from_parsed_slides(slides=slides)

def create_ppt_from_parsed_slides(slides,template_path = "Theme.pptx",  output_file="slide_deck.pptx"):
    prs = Presentation(template_path)

    for slide in slides:
        layout = prs.slide_layouts[0] if slide["is_title"] else prs.slide_layouts[1]
        ppt_slide = prs.slides.add_slide(layout)
        ppt_slide.shapes.title.text = slide["title"]

        # Add bullets
        if not slide["is_title"] and slide["bullets"]:
            body = ppt_slide.placeholders[1]
            body.text = "\n".join(slide["bullets"])

        # Add table
        if slide["table"]:
            rows = len(slide["table"])
            cols = len(slide["table"][0])
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(0.8)

            table = ppt_slide.shapes.add_table(rows, cols, left, top, width, height).table

            for i, row in enumerate(slide["table"]):
                for j, cell in enumerate(row):
                    table.cell(i, j).text = cell

    prs.save(output_file)
    print(f"Slide deck saved to: {output_file}")

def extract_table_data(text):
    """
    Extract table data from markdown-formatted text and convert to pandas DataFrame.
    """
    import pandas as pd
    
    # Find table lines (lines starting with |)
    table_lines = [line.strip() for line in text.split('\n') if line.strip().startswith('|')]
    
    if len(table_lines) < 2:  # Need at least header and one data row
        return None
    
    # Extract headers
    header_line = table_lines[0]
    headers = [cell.strip() for cell in header_line.split('|') if cell.strip()]
    
    # Skip separator line if present (line with dashes like | --- | --- |)
    data_start_idx = 1
    if len(table_lines) > 1 and all(cell.strip().replace('-', '').replace(':', '') == '' 
                                   for cell in table_lines[1].split('|') if cell.strip()):
        data_start_idx = 2
    
    # Extract data rows
    data = []
    for line in table_lines[data_start_idx:]:
        # Split by | and remove empty cells at the beginning/end
        cells = [cell.strip() for cell in line.split('|')]
        cells = [cell for cell in cells if cell]  # Remove empty strings
        
        if cells:
            # Make sure we have the same number of cells as headers
            if len(cells) < len(headers):
                cells.extend([''] * (len(headers) - len(cells)))
            elif len(cells) > len(headers):
                cells = cells[:len(headers)]
                
            data.append(cells)
    
    # Create DataFrame
    if data:
        df = pd.DataFrame(data, columns=headers)
        
        return df
    
    return None

#PPT Creator Agent
agent_e_chain = create_chain()

# State schema
class AgentState(TypedDict):
    question: str
    answer: Optional[str]
    ppt: Optional[bool]
    ppt_file_content: Optional[str]
    filename: Optional[str]
    excel : Optional[bool]
    excel_path : Optional[str]
    excel_name : Optional[str]
    agent_type : Optional[str]

# Define agent functions
def agent_a(state):
    result = agent_a_chain.invoke({"question": state["question"]})
    return {"answer": result["answer"], "agent_type" : "Answer from the ML Agent"}

def agent_b(state):
    result = agent_b_chain.invoke({"question": state["question"]})
    return {"answer": result["answer"], "agent_type" : "Answer from the Medical Agent"}

def agent_c(state):
    result = agent_c_chain.invoke({"question": state["question"]})
    return {"answer": result["answer"], "agent_type" : "Answer from the Document Agent"} 

def agent_d(state):
    """Finance agent that extracts data from finance PDF, formats it as a table, and saves it to Excel."""
    # Get the question
    question = state["question"]
    
    # Modify question to explicitly request tabular data
    table_question = f"{question} (FORMAT RESPONSE AS A MARKDOWN TABLE WITH COLUMN HEADERS)"
    
    # Get answer from the chain
    result = agent_d_chain({"question": table_question})
    answer = result["answer"]
    
    # Check if the answer has a table format (contains | characters and header separator)
    if '|' in answer and any(line.strip().startswith('|') and line.strip().endswith('|') for line in answer.split('\n')):
        try:
            # Extract the table data to pandas DataFrame
            df = extract_table_data(answer)
            
            if df is not None and not df.empty:
                # Save the DataFrame to Excel
                excel_path = 'finance_data.xlsx'
                df.to_excel(excel_path, index=False)
                print(f"Successfully saved data to {excel_path}")
                
                # Add a note to the answer about the Excel file
                answer += f"\n\nThe data has been saved to {excel_path}"
        except Exception as e:
            print(f"Error processing table data: {e}")
    else:
        print("Response did not contain proper table format")
    
    return {"answer": answer, "excel": True, "excel_path" : "finance_data.xlsx", "excel_name" : "finance_data.xlsx", "agent_type" : "Answer from the Finance Agent"}

def agent_e(state):
    result = agent_e_chain({"question": state["question"]})
    print("result",result)
    text = result["answer"]
    slide_content = extract_slide_content(text)
    print("Slide Content:", slide_content)
    # create_presentation(slide_content, filename="presentation.pptx")
    parse_slide_content(slide_content)
    response = {"answer": "The created slide deck is.", "ppt":True, "ppt_file_content":"slide_deck.pptx","filename":"presentation.pptx", "agent_type" : "Answer from the Slide Deck Agent"}
    print("Resp",response)
    return response

def agent_f(state):
    # Connect to SQLite database
    db = SQLDatabase.from_uri("sqlite:///./Chinook.db")
    llm = ChatGroq(
        groq_api_key=api_key,
        model_name= "llama-3.3-70b-versatile",
        temperature=0.4
    )
    
    # Create prompt for SQL generation
    prompt = ChatPromptTemplate.from_template(
        """
    You are a helpful assistant that translates natural language into SQL queries.
    Use the following database schema:
    {schema}
    
    Translate the following question into a SQL query:
    Question: {question}
    
    SQL Query:
    Return only the SQL query in response and nothing else.
    Don't add any backticks in response.
    """
    )
    chain = LLMChain(llm=llm, prompt=prompt)
    schema = db.get_table_info()
    sql_query = chain.run({"schema": schema, "question": state["question"]})
    print("Generated SQL Query:\n", sql_query)
    conn = sqlite3.connect('Chinook.db')
    df = pd.read_sql_query(sql_query, conn)
    conn.close()
    return {"answer": df.to_markdown(index=False), "agent_type" : "Answer from the SQL Agent"}

def error_handler_agent(state):
    return {"answer": "The question seems to be out of context, please ask something else."}

# Agent descriptions for the router
agent_descriptions = {
    "agent_a": "Expert in answering question about transformers. Is a research paper for Attention is all you need which talks about transformers.",
    "agent_b": "Expert in answering questions related to Radiotherapy (Medical).",
    "agent_c": "Expert in answering questions related to Software Development Life Cycle also known as SDLC.",
    "agent_d": "A financial data expert that extracts and presents financial information in table format and saves data to Excel. Handles queries related to financial reports, metrics, and data extraction.",
    "agent_e": """An expert in creating powerpoint presentations or slide decks related to I2E Consulting.
     Fetch all the necessary details to include in the slide deck. Try top explain each slide with multiple bulllet points to make the slide deck more interactive.""",
    "agent_f": "A SQL expert agent which is capable of answering questions related to Album, Artist, Customer, Employee, Genre, Invoice, InvoiceLine, MediaType, Playlist, PlaylistTrack, Track as these are the tables present in the Chinook database."
}

def supervisor_router(state):
    question = state["question"]
    system_prompt = """
    You are a router for a multi-agent system. Given a user question and a list of available agents with their descriptions,
    choose the best agent to handle the query.

    Return only the agent name from the list below.

    If the user question is not directly related to any of the agents, return \"No agent selected\".

    Agents:
    {agent_list}

    User question:
    {question}
    """
    prompt = PromptTemplate.from_template(system_prompt)
    formatted = prompt.format(
        agent_list="\n".join([f"{name}: {desc}" for name, desc in agent_descriptions.items()]),
        question=question
    )
    response = llm.invoke(formatted)
    selected = response.content.strip().lower()
    print("Understanding of the router agent:", selected)
    
    # Sanitize response to match node names
    if "agent_a" in selected:
        return "agent_a"
    elif "agent_b" in selected:
        return "agent_b"
    elif "agent_c" in selected:
        return "agent_c"
    elif "agent_d" in selected:
        return "agent_d"
    elif "agent_e" in selected:
        return "agent_e"
    elif "agent_f" in selected:
        return "agent_f"    
    else:
        print("Not a correct question")
        return "error_handler_agent"

# Build the multi-agent workflow
@st.cache_resource
def create_workflow():
    workflow = StateGraph(state_schema=AgentState)
    workflow.add_node("supervisor", supervisor_router)
    workflow.add_node("agent_a", agent_a)
    workflow.add_node("agent_b", agent_b)
    workflow.add_node("agent_c", agent_c)
    workflow.add_node("agent_d", agent_d)
    workflow.add_node("agent_e", agent_e)
    workflow.add_node("agent_f", agent_f)
    workflow.add_node("error_handler_agent", error_handler_agent)
    workflow.set_conditional_entry_point(supervisor_router)

    workflow.add_edge("agent_a", END)
    workflow.add_edge("agent_b", END)
    workflow.add_edge("agent_c", END)
    workflow.add_edge("agent_d", END)
    workflow.add_edge("agent_e", END)
    workflow.add_edge("agent_f", END)
    workflow.add_edge("error_handler_agent", END)
    return workflow.compile()

graph = create_workflow()

# # Streamlit UI
# st.title("Multi-Agent QA System")
# st.write("Ask a question about transformers (ML), radiotherapy, software development, or the Chinook database.")

# query = st.text_input("Enter your question:")
# if st.button("Get Answer"):
#     if query:
#         with st.spinner("Processing your query..."):
#             result = graph.invoke({"question": query})
#             answer = result.get("answer", "No answer returned.")
#         st.subheader("Answer:")
#         st.write(answer)
#     else:
#         st.warning("Please enter a question before clicking the button.")
# ─── STREAMLIT UI (with agent_type, PPT & Excel logic) ───
st.title("Multi-Agent QA System")
st.write("Ask a question about transformers, radiotherapy, software development, or the Chinook database.")

query = st.text_input("Enter your question:")
if st.button("Get Answer"):
    if not query.strip():
        st.warning("Please enter a question before clicking the button.")
    else:
        # Directly invoke the graph (no requests)
        result = graph.invoke({"question": query})

        # 1) Show which agent answered (agent_type)
        #    We assume that your agent functions now return something like:
        #    {
        #      "answer": "…",
        #      "agent_type": "🏥 [Radiotherapy Expert]",
        #      "ppt": True or False,
        #      "ppt_file_content": "/path/to/generated.pptx",
        #      "filename": "my_presentation.pptx",
        #      "excel": True or False,
        #      "excel_path": "/path/to/generated.xlsx",
        #      "excel_name": "report.xlsx"
        #    }
        agent_type = result.get("agent_type", None)
        if agent_type:
            st.subheader(agent_type)

        # 2) Always show the textual answer
        answer = result.get("answer", "No answer returned.")
        st.write(answer)

        # 3) If your agents have generated a PPT, show a download button
        if result.get("ppt", False):
            ppt_path = result.get("ppt_file_content")
            ppt_filename = result.get("filename", "presentation.pptx")
            try:
                with open(ppt_path, "rb") as f:
                    ppt_bytes = f.read()
                st.download_button(
                    label="Download PPT",
                    data=ppt_bytes,
                    file_name=ppt_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"Could not open PPT at {ppt_path}: {e}")

        # 4) If your agents have generated an Excel file, show a download button
        if result.get("excel", False):
            excel_path = result.get("excel_path")
            excel_filename = result.get("excel_name", "data.xlsx")
            try:
                with open(excel_path, "rb") as f:
                    excel_bytes = f.read()
                st.download_button(
                    label="Download Excel",
                    data=excel_bytes,
                    file_name=excel_filename,
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )
            except Exception as e:
                st.error(f"Could not open Excel at {excel_path}: {e}")
